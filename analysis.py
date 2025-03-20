import streamlit as st
import pandas as pd
import numpy as np
from scipy.stats import chi2_contingency, fisher_exact, ttest_ind, mannwhitneyu
from statsmodels.stats.proportion import proportions_ztest
from sklearn.model_selection import train_test_split
from sklearn.metrics import accuracy_score, classification_report
from xgboost import XGBClassifier
import os
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
import requests
import tempfile

st.set_page_config(page_title="A/B Testing Analysis")

st.title("A/B Testing Data Analysis")
st.write(
    "Upload an A/B testing dataset, map columns, perform statistical tests, "
    "train an ML model, and generate reports using your PowerPoint template."
)

def show_column_description():
    """Displays dataset column descriptions."""
    with st.expander("View Dataset Guidelines and Use Cases", expanded=False):
        st.subheader("Dataset Structure & A/B Testing Analysis Features")
        st.markdown("""
            **Required Columns:** `user_id`, `group`, `engagement`, `conversion`, `metric`  
            **Support for Extra Columns:** Other columns can be used in the ML model  
            **Dynamic Metric Selection:** Users can choose **any column** as the main metric  
            **Machine Learning with All Features:** Other columns can be used in ML analysis  
            **PowerPoint Report Update:** Uses a **template-based approach** for structured reporting  
            """)

        column_data = {
            "Column Name": ["user_id", "group", "engagement", "conversion", "metric"],
            "Data Type": ["Integer", "String", "Binary (0/1)", "Binary (0/1)", "Float"],
            "Description": [
                "Unique identifier for each user",
                "Control (A) or Treatment (B)",
                "Whether the user interacted with the experiment (Clicked, Opened, etc.)",
                "Whether the user completed the desired action (Purchased, Signed Up, etc.)",
                "The key performance indicator (KPI) for the experiment (e.g., time spent, clicks, revenue, session duration, etc.)"
            ]
        }

        st.table(pd.DataFrame(column_data).set_index("Column Name"))

def upload_and_map_columns():
    """Handles file upload, column mapping, and feature selection."""
    st.subheader("Upload A/B Testing Data")
    uploaded_file = st.file_uploader("Upload CSV", type=["csv"])

    if uploaded_file is None:
        return None, None

    df = pd.read_csv(uploaded_file)
    st.write("Preview of Uploaded Data")
    st.dataframe(df, hide_index=True)

    required_columns = ["user_id", "group", "engagement", "conversion"]
    # Allow user to select **one main metric column**
    metric_column = st.selectbox("Select the main metric column", df.columns)
    
    column_mapping = {}
    for col in required_columns:
        column_mapping[col] = st.selectbox(f"Select column for {col}", df.columns, key=col)

    df = df.rename(columns=column_mapping)
    st.session_state["selected_metric_name"] = metric_column
    df.rename(columns={metric_column: "metric"}, inplace=True)

    missing_cols = [col for col in required_columns if col not in df.columns]
    if missing_cols:
        st.error(f"Missing columns: {missing_cols}. Please map them correctly.")
        return None, None

    df["group"] = df["group"].astype(str)
    if not set(df["group"].unique()).issubset({"A", "B"}):
        st.error("Group column should only contain 'A' and 'B'. Please correct your data.")
        return None, None

    df["engagement"] = df["engagement"].astype(int)
    df["conversion"] = df["conversion"].astype(int)
    df["metric"] = df["metric"].astype(float)

    st.success("Data uploaded and validated.")

    # Feature selection
    all_columns = list(df.columns)
    default_features = ["group", "engagement", "metric"]
    selected_features = st.multiselect("Select additional columns for ML:", all_columns, default=default_features)

    return df, selected_features

def compute_summary_statistics(df):
    """Computes conversion rates for A/B test groups."""
    
    summary = df.groupby("group")["conversion"].agg(['sum', 'count'])
    summary['conversion_rate'] = summary['sum'] / summary['count']
    return summary


def perform_statistical_test(df, test_choice):
    """Executes the selected statistical test and returns the p-value."""
    
    summary = compute_summary_statistics(df)

    if test_choice == "Chi-Square Test":
        contingency_table = pd.crosstab(df["group"], df["conversion"])
        _, p_value, _, _ = chi2_contingency(contingency_table)
    elif test_choice == "Fisher’s Exact Test":
        contingency_table = pd.crosstab(df["group"], df["conversion"])
        _, p_value = fisher_exact(contingency_table)
    elif test_choice == "Z-Test for Proportions":
        count = np.array([summary.loc["A", "sum"], summary.loc["B", "sum"]])
        nobs = np.array([summary.loc["A", "count"], summary.loc["B", "count"]])
        _, p_value = proportions_ztest(count, nobs, alternative='two-sided')
    elif test_choice == "T-Test":
        _, p_value = ttest_ind(
            df[df["group"] == "A"]["metric"], 
            df[df["group"] == "B"]["metric"], 
            equal_var=False
        )
    elif test_choice == "Mann-Whitney U Test":
        _, p_value = mannwhitneyu(
            df[df["group"] == "A"]["metric"], 
            df[df["group"] == "B"]["metric"], 
            alternative='two-sided'
        )
    else:
        p_value = None

    return p_value


def prepare_data(df, selected_features):
    """Prepares data for ML with dynamic feature selection and encoding."""
    
    df["group"] = df["group"].map({"A": 0, "B": 1})

    # Identify categorical features
    categorical_features = [col for col in selected_features if df[col].dtype == "object"]

    # Apply one-hot encoding
    df = pd.get_dummies(df, columns=categorical_features, drop_first=True)

    # Find the new encoded column names
    updated_feature_names = df.columns.tolist()

    # Adjust selected features to match the new column names
    selected_features = [
        feature for feature in updated_feature_names if any(orig in feature for orig in selected_features)
    ]

    # Check if any valid feature remains
    if not selected_features:
        st.error("No valid features selected for training. Please select at least one column.")
        return None, None

    X = df[selected_features]
    y = df["conversion"]

    return X, y


def train_xgboost_model(X_train, y_train):
    """Trains an XGBoost classification model."""
    model = XGBClassifier(
        objective="binary:logistic",
        eval_metric="logloss",  
        n_estimators=100,
        learning_rate=0.05,
        max_depth=5,
        random_state=42
    )
    model.fit(X_train, y_train)
    return model

def perform_z_test_for_proportions(summary):
    """
    Performs a Z-test for comparing conversion rates between A/B groups.

    :param summary: Summary DataFrame containing conversion counts
    :return: p-value of the Z-Test
    """
    # Ensure correct index labels ("A"/"B" or 0/1)
    group_A = "A" if "A" in summary.index else 0
    group_B = "B" if "B" in summary.index else 1  

    # Extract counts and total observations
    count = np.array([summary.loc[group_A, "sum"], summary.loc[group_B, "sum"]])
    nobs = np.array([summary.loc[group_A, "count"], summary.loc[group_B, "count"]])

    # Perform Z-test
    _, p_value = proportions_ztest(count, nobs, alternative="two-sided")

    return p_value

def evaluate_xgboost_model(model, X_test, y_test):
    """Evaluates the trained XGBoost model and stores results in session state."""
    
    y_pred = model.predict(X_test)
    accuracy = accuracy_score(y_test, y_pred)
    
    # Generate classification report
    report = classification_report(y_test, y_pred, output_dict=True, zero_division=0)
    report_df = pd.DataFrame(report).transpose()

    # Get actual feature names from X_test
    feature_names = X_test.columns.tolist()

    # Ensure feature importance matches the number of features
    if len(feature_names) != len(model.feature_importances_):
        st.error("Feature importance calculation error: Feature names do not match model inputs.")
        return
    
    # Feature Importance
    importance_df = pd.DataFrame({"Feature": feature_names, "Importance": model.feature_importances_})
    importance_df = importance_df.sort_values(by="Importance", ascending=False)

    # Store results in session state
    st.session_state["accuracy"] = accuracy
    st.session_state["classification_report"] = report_df
    st.session_state["feature_importance"] = importance_df

def update_powerpoint_template(template_path, summary, model):
    """Updates an existing PowerPoint template with A/B test results and ML insights."""
    
    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text

                # Update Conversion Rates
                if "Conversion Rate (Group A):" in text:
                    shape.text_frame.text = (
                        f"• Conversion Rate (A): {summary.loc['A', 'conversion_rate']:.2%}\n"
                        f"• Conversion Rate (B): {summary.loc['B', 'conversion_rate']:.2%}"
                    )

                # Update Model Accuracy
                if "Model Accuracy:" in text and "accuracy" in st.session_state:
                    shape.text_frame.text = f"Model Accuracy: {st.session_state['accuracy']:.2%}"

    # **Step 2: Generate Updated Charts**
    conversion_chart = "conversion_rate_chart.png"
    plt.figure(figsize=(6, 4))
    sns.barplot(x=summary.index, y=summary["conversion_rate"], palette=["blue", "green"])
    plt.title("Conversion Rate by Group")
    plt.ylabel("Conversion Rate")
    plt.xlabel("Group (A = Control, B = Treatment)")
    plt.ylim(0, 1)
    plt.savefig(conversion_chart, bbox_inches="tight", dpi=300)
    plt.close()

    # Ensure feature importance data exists
    if "feature_importance" not in st.session_state or st.session_state["feature_importance"] is None:
        st.error("Feature importance data not found. Train the model first.")
        return None

    # Feature Importance Chart
    importance_df = st.session_state["feature_importance"]

    plt.figure(figsize=(6, 4))
    sns.barplot(x=importance_df["Importance"], y=importance_df["Feature"], palette="Blues_r")
    plt.title("Feature Importance")
    plt.xlabel("Importance Score")
    plt.ylabel("Feature")

    image_path = "feature_importance.png"
    plt.savefig(image_path, bbox_inches="tight", dpi=300)
    plt.close()

    # Replace images in the template
    image_replaced = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Shape type 13 is a picture
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                slide.shapes._spTree.remove(shape._element)
                slide.shapes.add_picture(image_path, left, top, width=width, height=height)
                image_replaced = True
                break  # Replace only the first image found

    # Ensure at least one image was replaced
    if not image_replaced:
        st.warning("No image found to replace in the PowerPoint template.")

    # Save updated PPT
    ppt_filename = "Updated_AB_Test_Report.pptx"
    prs.save(ppt_filename)
    
    # Remove temp image
    os.remove(image_path)

    return ppt_filename

def generate_powerpoint_template(template_path, summary, test_results, model, X_test):
    """
    Updates an existing PowerPoint template with new A/B test results and ML model insights.

    :param template_path: Path to the PowerPoint template.
    :param summary: DataFrame with conversion rates.
    :param test_results: Dictionary with statistical test p-values.
    :param model: Trained ML model.
    :param X_test: Test dataset features.
    :return: Updated PowerPoint file path.
    """
    prs = Presentation(template_path)

    # **Step 1: Update Text in Placeholders**
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text

                # Update A/B Test Summary
                if "Conversion Rate (Group A):" in text:
                    shape.text_frame.text = (
                        f"• Conversion Rate (Group A): {summary.loc['A', 'conversion_rate']:.2%}\n"
                        f"• Conversion Rate (Group B): {summary.loc['B', 'conversion_rate']:.2%}"
                    )

                # Update Statistical Test Results
                if "P-Values from Statistical Tests:" in text:
                    updated_text = "P-Values from Statistical Tests:\n"
                    for test, p_value in test_results.items():
                        significance = "(Significant)" if p_value < 0.05 else "(Not Significant)"
                        updated_text += f"{test}: {p_value:.4f} {significance}\n"
                    shape.text_frame.text = updated_text

    # **Step 2: Generate Updated Charts**
    conversion_chart = "conversion_rate_chart.png"
    plt.figure(figsize=(6, 4))
    sns.barplot(x=summary.index, y=summary["conversion_rate"], palette=["blue", "green"])
    plt.title("Conversion Rate by Group")
    plt.ylabel("Conversion Rate")
    plt.xlabel("Group (A = Control, B = Treatment)")
    plt.ylim(0, 1)
    plt.savefig(conversion_chart, bbox_inches="tight", dpi=300)
    plt.close()

    feature_importance_chart = "feature_importance_chart.png"
    importance_df = pd.DataFrame({"Feature": X_test.columns, "Importance": model.feature_importances_})
    plt.figure(figsize=(6, 4))
    sns.barplot(x=importance_df["Importance"], y=importance_df["Feature"], palette="Blues_r")
    plt.title("Feature Importance for Conversion Prediction")
    plt.xlabel("Importance Score")
    plt.ylabel("Feature")
    plt.savefig(feature_importance_chart, bbox_inches="tight", dpi=300)
    plt.close()

    # **Step 3: Replace Images in the Template**
    image_count = 0  # Counter to track which image to replace
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # If it's a picture
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                slide.shapes._spTree.remove(shape._element)  # Remove the existing image

                # Add new image in the same position
                if image_count == 0:
                    slide.shapes.add_picture(conversion_chart, left, top, width=width, height=height)
                    image_count += 1
                elif image_count == 1:
                    slide.shapes.add_picture(feature_importance_chart, left, top, width=width, height=height)
                    image_count += 1

    # **Step 4: Save the Updated Presentation**
    updated_ppt_path = "Updated_AB_Test_Report.pptx"
    prs.save(updated_ppt_path)

    # Remove temporary images
    os.remove(conversion_chart)
    os.remove(feature_importance_chart)

    return updated_ppt_path


def main():
    """Streamlit app workflow."""
    for key in ["model", "X_test", "y_test", "summary", "accuracy", "classification_report", "feature_importance"]:
        if key not in st.session_state:
            st.session_state[key] = None

    show_column_description()

    df, selected_features = upload_and_map_columns()



    if df is not None and selected_features:
        st.subheader("Data Ready for Analysis")

        summary = compute_summary_statistics(df)
        st.write("Conversion Rate Summary")
        st.table(summary)

        # Statistical Test Selection
        test_choice = st.radio(
            "Select a Statistical Test",
            [
                "Chi-Square Test",
                "Fisher’s Exact Test",
                "T-Test",
                "Mann-Whitney U Test",
            ]
        )
        # Perform selected test
        p_value = (
            perform_z_test_for_proportions(summary) if test_choice == "Z-Test for Proportions" 
            else perform_statistical_test(df, test_choice)
        )

        if p_value is not None:
            st.write("Test Result")
            st.write(f"P-Value: {p_value:.4f}")

            if p_value < 0.05:
                st.success("The result is statistically significant (p < 0.05).")
            else:
                st.warning("The result is not statistically significant (p ≥ 0.05).")
        t = "Z-Test for Proportions"
        p_value = (
            perform_z_test_for_proportions(summary) if t == "Z-Test for Proportions" 
            else perform_statistical_test(df, t)
        )

        if p_value is not None:
            st.write("Z-Test for Proportions Result")
            st.write(f"P-Value: {p_value:.4f}")

            if p_value < 0.05:
                st.success("The result is statistically significant (p < 0.05).")
            else:
                st.warning("The result is not statistically significant (p ≥ 0.05).")

        if st.session_state["model"] is None:
            if st.button("Train ML Model"):
                X, y = prepare_data(df, selected_features)
                X_train, X_test, y_train, y_test = train_test_split(X, y, test_size=0.2, random_state=42, stratify=y)
                model = train_xgboost_model(X_train, y_train)
                evaluate_xgboost_model(model, X_test, y_test)
                st.session_state["model"] = model
                st.session_state["summary"] = summary
                st.session_state["X_test"] = X_test
                st.session_state["y_test"] = y_test
                st.rerun()

        if st.session_state["model"] is not None:
            st.subheader("Model Evaluation Results")

            if st.session_state["accuracy"] is not None:
                st.write(f"**Model Accuracy:** {st.session_state['accuracy']:.4f}")
            else:
                st.warning("Model accuracy is not available yet. Train the model first.")

            if st.session_state["classification_report"] is not None:
                st.write("**Classification Report**")
                st.dataframe(st.session_state["classification_report"])

            if st.session_state["feature_importance"] is not None:
                st.write("**Feature Importance**")
                st.write(f"**Selected Metric for Analysis:** `{st.session_state['selected_metric_name']}`")
                st.bar_chart(st.session_state["feature_importance"].set_index("Feature"))

            if st.button("Generate PowerPoint Report"):
                TEMPLATE_URL = "https://raw.githubusercontent.com/pspreethi/A-B-testing-analysis-streamlit-app/main/template.pptx"

                # Download file and save it temporarily
                response = requests.get(TEMPLATE_URL)
                if response.status_code == 200:
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                        tmp.write(response.content)
                        template_path = tmp.name  # Use this path
                    prs = Presentation(template_path)
                    print("Template loaded successfully!")
                else:
                    print("Failed to download the template.")

                test_results = {
                    f"{test_choice} Test": perform_statistical_test(df, test_choice),
                    "Z-Test for Proportions": perform_z_test_for_proportions(summary)
                }

                ppt_filename = generate_powerpoint_template(template_path, summary, test_results, st.session_state["model"], st.session_state["X_test"])

                with open(ppt_filename, "rb") as ppt_file:
                    st.download_button("Download Report", data=ppt_file, file_name=ppt_filename, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")

if __name__ == "__main__":
    main()
